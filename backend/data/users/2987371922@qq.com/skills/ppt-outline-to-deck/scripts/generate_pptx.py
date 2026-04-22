#!/usr/bin/env python3
"""Build editable PPTX from deck.json and slide images."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def _load_deck(path: Path) -> dict[str, Any]:
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        raise RuntimeError(f"无法读取 deck_json: {exc}") from exc
    if not isinstance(raw, dict):
        raise RuntimeError("deck_json 必须是 JSON 对象")
    slides = raw.get("slides")
    if not isinstance(slides, list) or not slides:
        raise RuntimeError("deck_json 缺少非空 slides[]")
    return raw


def _safe_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    return str(value)


def build_pptx(deck: dict[str, Any], out_file: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError(
            "缺少 python-pptx 依赖，请在用户 sandbox requirements 中安装 python-pptx"
        ) from exc

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank

    deck_meta = deck.get("deck_meta") if isinstance(deck.get("deck_meta"), dict) else {}
    deck_title = _safe_text(deck_meta.get("title") or "PPT")

    # Cover slide
    cover = prs.slides.add_slide(slide_layout)
    t_box = cover.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.0), Inches(1.0))
    tf = t_box.text_frame
    tf.text = deck_title
    tf.paragraphs[0].font.size = Pt(38)
    tf.paragraphs[0].font.bold = True

    subtitle = _safe_text(deck_meta.get("subtitle"))
    if subtitle:
        s_box = cover.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.0), Inches(1.2))
        sf = s_box.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(20)

    for slide_data in deck.get("slides", []):
        if not isinstance(slide_data, dict):
            continue
        slide = prs.slides.add_slide(slide_layout)

        title = _safe_text(slide_data.get("title"))
        bullets = slide_data.get("bullets")
        notes = _safe_text(slide_data.get("speaker_notes"))
        image_path = _safe_text(slide_data.get("image_path"))

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.size = Pt(30)

        left = Inches(0.6)
        top = Inches(1.3)
        width = Inches(5.2)
        height = Inches(4.8)
        if image_path:
            image_file = Path(image_path)
            if not image_file.is_absolute():
                image_file = Path.cwd() / image_file
            if image_file.exists():
                slide.shapes.add_picture(str(image_file), Inches(6.0), Inches(1.3), width=Inches(6.8), height=Inches(4.2))
                width = Inches(5.0)

        bullet_box = slide.shapes.add_textbox(left, top, width, height)
        bullet_tf = bullet_box.text_frame
        bullet_tf.clear()
        if isinstance(bullets, list) and bullets:
            first = bullet_tf.paragraphs[0]
            first.text = _safe_text(bullets[0])
            first.font.size = Pt(20)
            for item in bullets[1:]:
                p = bullet_tf.add_paragraph()
                p.text = _safe_text(item)
                p.font.size = Pt(18)
                p.level = 0
        else:
            bullet_tf.text = ""

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    out_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_file))


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate editable PPTX from deck.json")
    parser.add_argument("--deck_json", required=True, help="Path to deck.json")
    parser.add_argument("--output", default="final.pptx", help="Output pptx path")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    deck_path = Path(args.deck_json).resolve()
    out_path = Path(args.output).resolve()
    try:
        deck = _load_deck(deck_path)
        build_pptx(deck, out_path)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    print(f"PPTX 已生成: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
